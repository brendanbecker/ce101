#!/usr/bin/env python3
"""
Analyze PowerPoint presentations for emoji and unicode issues.
"""

import sys
import unicodedata
from pptx import Presentation


def is_emoji(char):
    """Check if a character is an emoji."""
    # Common emoji unicode ranges
    code_point = ord(char)

    # Emoticons
    if 0x1F600 <= code_point <= 0x1F64F:
        return True
    # Miscellaneous Symbols and Pictographs
    if 0x1F300 <= code_point <= 0x1F5FF:
        return True
    # Transport and Map Symbols
    if 0x1F680 <= code_point <= 0x1F6FF:
        return True
    # Supplemental Symbols and Pictographs
    if 0x1F900 <= code_point <= 0x1F9FF:
        return True
    # Symbols and Pictographs Extended-A
    if 0x1FA70 <= code_point <= 0x1FAFF:
        return True
    # Additional emoji ranges
    if 0x2600 <= code_point <= 0x26FF:
        return True
    if 0x2700 <= code_point <= 0x27BF:
        return True
    # Enclosed characters that might be emoji
    if 0x24C2 <= code_point <= 0x1F251:
        return True

    return False


def is_broken_unicode(char):
    """Check if a character appears to be broken/garbled unicode."""
    # Replacement character
    if char == '\ufffd':
        return True

    # Check for private use areas (often indicate broken encoding)
    code_point = ord(char)
    if 0xE000 <= code_point <= 0xF8FF:  # Private Use Area
        return True
    if 0xF0000 <= code_point <= 0xFFFFD:  # Supplementary Private Use Area-A
        return True
    if 0x100000 <= code_point <= 0x10FFFD:  # Supplementary Private Use Area-B
        return True

    # Check if character category is "Other" which might indicate issues
    category = unicodedata.category(char)
    if category in ['Co', 'Cn']:  # Private use or not assigned
        return True

    return False


def get_char_info(char):
    """Get detailed information about a character."""
    code_point = ord(char)
    try:
        name = unicodedata.name(char)
    except ValueError:
        name = "UNNAMED"

    category = unicodedata.category(char)

    return {
        'char': char,
        'code_point': f'U+{code_point:04X}',
        'name': name,
        'category': category,
        'is_emoji': is_emoji(char),
        'is_broken': is_broken_unicode(char)
    }


def extract_text_from_shape(shape):
    """Extract text from a shape, handling different types."""
    text_parts = []

    if hasattr(shape, "text"):
        text_parts.append(shape.text)

    if hasattr(shape, "text_frame"):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_parts.append(run.text)

    return '\n'.join(text_parts)


def analyze_presentation(pptx_path):
    """Analyze a PowerPoint presentation for emoji and unicode issues."""
    print(f"\n{'='*80}")
    print(f"Analyzing: {pptx_path}")
    print(f"{'='*80}\n")

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"ERROR: Could not open presentation: {e}")
        return

    findings = []
    total_emoji_count = 0
    total_broken_count = 0

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_findings = {
            'slide_num': slide_idx,
            'emoji_chars': [],
            'broken_chars': [],
            'other_special': []
        }

        # Extract text from all shapes in the slide
        for shape in slide.shapes:
            text = extract_text_from_shape(shape)

            if not text:
                continue

            # Analyze each character
            for char in text:
                # Skip common ASCII and whitespace
                if ord(char) < 128 or char.isspace():
                    continue

                char_info = get_char_info(char)

                if char_info['is_broken']:
                    slide_findings['broken_chars'].append(char_info)
                    total_broken_count += 1
                elif char_info['is_emoji']:
                    slide_findings['emoji_chars'].append(char_info)
                    total_emoji_count += 1
                else:
                    # Check for other potentially interesting unicode
                    if ord(char) > 0x2000:  # Beyond common Latin/symbols
                        slide_findings['other_special'].append(char_info)

        # Remove duplicates
        slide_findings['emoji_chars'] = list({c['code_point']: c for c in slide_findings['emoji_chars']}.values())
        slide_findings['broken_chars'] = list({c['code_point']: c for c in slide_findings['broken_chars']}.values())
        slide_findings['other_special'] = list({c['code_point']: c for c in slide_findings['other_special']}.values())

        if slide_findings['emoji_chars'] or slide_findings['broken_chars'] or slide_findings['other_special']:
            findings.append(slide_findings)

    # Print results
    print(f"Total Slides: {len(prs.slides)}")
    print(f"Slides with Issues: {len(findings)}")
    print(f"Total Emoji Characters Found: {total_emoji_count}")
    print(f"Total Broken Characters Found: {total_broken_count}")
    print()

    if not findings:
        print("No emoji or unicode issues found!")
        return

    for finding in findings:
        print(f"\n--- Slide {finding['slide_num']} ---")

        if finding['broken_chars']:
            print(f"\n  BROKEN/GARBLED UNICODE ({len(finding['broken_chars'])} unique):")
            for char_info in finding['broken_chars']:
                print(f"    {char_info['char']!r} - {char_info['code_point']} - {char_info['name']} ({char_info['category']})")

        if finding['emoji_chars']:
            print(f"\n  EMOJI CHARACTERS ({len(finding['emoji_chars'])} unique):")
            for char_info in finding['emoji_chars']:
                print(f"    {char_info['char']} - {char_info['code_point']} - {char_info['name']}")

        if finding['other_special']:
            print(f"\n  OTHER SPECIAL UNICODE ({len(finding['other_special'])} unique):")
            for char_info in finding['other_special']:
                print(f"    {char_info['char']} - {char_info['code_point']} - {char_info['name']}")

    print(f"\n{'='*80}\n")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python analyze_pptx_unicode.py <presentation.pptx> [presentation2.pptx ...]")
        sys.exit(1)

    for pptx_file in sys.argv[1:]:
        analyze_presentation(pptx_file)
