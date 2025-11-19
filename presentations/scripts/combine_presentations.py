#!/usr/bin/env python3
"""
Combine 7 CE101 PowerPoint presentations into one master presentation.
Uses a more robust approach by directly manipulating presentation parts.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from copy import deepcopy
import os

def create_section_divider(prs, module_number, module_title):
    """
    Create a section divider slide with module number and title.
    Uses the first available layout as a base.
    """
    # Use the first available layout
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Clear all shapes from the slide
    for shape in list(slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Set background color to a dark blue/gray
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(31, 56, 100)  # Dark blue
    except:
        pass  # If background setting fails, continue anyway

    # Add module text
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # Add module number
    p1 = text_frame.paragraphs[0]
    p1.text = f"Module {module_number}"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.alignment = PP_ALIGN.CENTER

    # Add module title
    p2 = text_frame.add_paragraph()
    p2.text = module_title
    p2.font.size = Pt(54)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(255, 255, 255)
    p2.alignment = PP_ALIGN.CENTER

    return slide

def clone_slide(prs, source_slide):
    """
    Clone a slide by copying its XML structure.
    This preserves all formatting, shapes, and design elements.
    """
    from lxml import etree

    # Get the source slide's XML
    source_slide_part = source_slide.part

    # Add a blank slide to destination
    blank_layout = prs.slide_layouts[0]
    dest_slide = prs.slides.add_slide(blank_layout)
    dest_slide_part = dest_slide.part

    # Copy the slide's shape tree
    source_tree = source_slide.element.cSld.spTree
    dest_tree = dest_slide.element.cSld.spTree

    # Clear destination shape tree
    for sp in list(dest_tree):
        dest_tree.remove(sp)

    # Copy all shapes from source
    for sp in source_tree:
        dest_tree.append(deepcopy(sp))

    # Copy background if present
    try:
        if hasattr(source_slide.element.cSld, 'bg'):
            if hasattr(dest_slide.element.cSld, 'bg'):
                dest_slide.element.cSld.remove(dest_slide.element.cSld.bg)
            dest_slide.element.cSld.insert(0, deepcopy(source_slide.element.cSld.bg))
    except:
        pass

    return dest_slide

def main():
    slides_dir = "slides"

    # Source presentations in order
    sources = [
        "CE101-Session1-CoreConcepts.pptx",
        "CE101_Session2_Filesystem_Organization.pptx",
        "CE101_Session3_Multi_Tab_Orchestration.pptx",
        "CE101_Session4_Local_Data_Stores.pptx",
        "CE101_Session5_Integration_Patterns.pptx",
        "CE101_Session6_Practical_Patterns.pptx",
        "CE101_Session7_Common_Pitfalls.pptx"
    ]

    # Module information for section dividers
    modules = [
        None,  # Module 1 doesn't need a divider (it starts the deck)
        {"number": 2, "title": "Filesystem Organization"},
        {"number": 3, "title": "Multi-Tab Orchestration"},
        {"number": 4, "title": "Local Data Stores"},
        {"number": 5, "title": "Integration Patterns"},
        {"number": 6, "title": "Practical Patterns"},
        {"number": 7, "title": "Common Pitfalls"}
    ]

    # Slide ranges to copy from each session (0-indexed)
    # Format: (start_index, end_index) - end_index is exclusive
    slide_ranges = [
        (0, 11),   # Session 1: Slides 1-11 (skip slide 12 - "Questions/Next")
        (1, 11),   # Session 2: Slides 2-11 (skip slide 1 - title, skip slide 12 - transition)
        (1, 12),   # Session 3: Slides 2-12 (skip slide 1 - title, skip slide 13 - transition)
        (1, 11),   # Session 4: Slides 2-11 (skip slide 1 - title, skip slide 12 - transition)
        (1, 15),   # Session 5: Slides 2-15 (skip slide 1 - title, skip slide 16 - transition)
        (1, 16),   # Session 6: Slides 2-16 (skip slide 1 - title, skip slide 17 - transition)
        (1, 20)    # Session 7: Slides 2-20 (skip slide 1 - title, keep journey/principles/thank you)
    ]

    print("Starting presentation combination...")
    print(f"Processing {len(sources)} presentations")

    # Load first presentation as template
    first_prs_path = os.path.join(slides_dir, sources[0])
    print(f"\nLoading template from: {first_prs_path}")

    # Create new presentation from first source
    combined_prs = Presentation(first_prs_path)

    # Remove all slides from the combined presentation
    xml_slides = combined_prs.slides._sldIdLst
    slides_to_remove = list(range(len(xml_slides) - 1, -1, -1))
    for i in slides_to_remove:
        rId = xml_slides[i].rId
        combined_prs.part.drop_rel(rId)
        del xml_slides[i]

    print(f"Created empty presentation template")

    # Process each session
    total_slides_added = 0

    for session_idx, source_file in enumerate(sources):
        session_num = session_idx + 1
        source_path = os.path.join(slides_dir, source_file)

        print(f"\n{'='*60}")
        print(f"Processing Session {session_num}: {source_file}")
        print(f"{'='*60}")

        # Load source presentation
        prs_source = Presentation(source_path)
        total_source_slides = len(prs_source.slides)
        print(f"Source has {total_source_slides} slides")

        # Add section divider for sessions 2-7
        if modules[session_idx] is not None:
            print(f"Adding section divider: Module {modules[session_idx]['number']}: {modules[session_idx]['title']}")
            try:
                create_section_divider(combined_prs, modules[session_idx]['number'], modules[session_idx]['title'])
                total_slides_added += 1
                print(f"  ✓ Section divider added")
            except Exception as e:
                print(f"  ✗ Error creating section divider: {e}")

        # Copy slides from this session
        start_idx, end_idx = slide_ranges[session_idx]
        slides_to_copy = list(range(start_idx, end_idx))

        print(f"Copying slides {start_idx + 1} to {end_idx} ({len(slides_to_copy)} slides)")

        for slide_idx in slides_to_copy:
            try:
                source_slide = prs_source.slides[slide_idx]
                clone_slide(combined_prs, source_slide)
                total_slides_added += 1
                print(f"  ✓ Copied slide {slide_idx + 1}")
            except Exception as e:
                print(f"  ✗ Error copying slide {slide_idx + 1}: {e}")
                import traceback
                traceback.print_exc()

    # Save combined presentation
    output_path = os.path.join(slides_dir, "CE101-Complete-Course.pptx")
    print(f"\n{'='*60}")
    print(f"Saving combined presentation to: {output_path}")
    print(f"Total slides in output: {total_slides_added}")
    print(f"{'='*60}")

    combined_prs.save(output_path)

    print(f"\n✓ Successfully created: {output_path}")
    print(f"Final slide count: {len(combined_prs.slides)}")

if __name__ == "__main__":
    main()
