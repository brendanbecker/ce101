#!/usr/bin/env python3
"""Extract text from PowerPoint files"""
from pptx import Presentation
import sys

if len(sys.argv) < 2:
    print("Usage: read_pptx.py <file.pptx>")
    sys.exit(1)

prs = Presentation(sys.argv[1])
print(f"Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"{'='*60}")
    print(f"SLIDE {i}")
    print(f"{'='*60}")

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text)
            print()
    print()
